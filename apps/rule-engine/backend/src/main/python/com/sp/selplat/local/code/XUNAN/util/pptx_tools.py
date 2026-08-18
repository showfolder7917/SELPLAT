#!/usr/bin/env python3
"""当前用户 PPTX 能力共享工具。

统一封装尺寸换算、可编辑文本、透明卡片、图片裁切、链接和预览渲染，
让各课程生成器只保留业务布局与内容编排。
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import shutil
import subprocess
import tempfile
from typing import Iterable

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


PIXELS_PER_INCH = 96


def px(value: float) -> int:
    """把现有生成器的 96 DPI 像素坐标转换为 PowerPoint EMU。"""

    return int(Inches(value / PIXELS_PER_INCH))


def rgb(value: str) -> RGBColor:
    """解析 `#RRGGBB` 或 `RRGGBB` 颜色。"""

    clean = value.lstrip("#")[:6]
    return RGBColor.from_string(clean.upper())


def create_presentation(width: int = 1280, height: int = 720) -> Presentation:
    """创建使用像素布局语义的空白演示文稿。"""

    presentation = Presentation()
    presentation.slide_width = px(width)
    presentation.slide_height = px(height)
    return presentation


def add_text(
    slide,
    text: str,
    box: tuple[float, float, float, float],
    *,
    font_size: float = 24,
    color: str = "#24363E",
    font_name: str = "Noto Sans CJK SC",
    bold: bool = False,
    align: str = "left",
    vertical: str = "middle",
    fill: str | None = None,
    transparency: int = 0,
    radius: bool = False,
    name: str | None = None,
):
    """添加可编辑文本框，并返回形状供调用方设置跳转。"""

    left, top, width, height = (px(item) for item in box)
    if fill:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        _set_fill_transparency(shape, transparency)
        shape.line.fill.background()
        frame = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(left, top, width, height)
        frame = shape.text_frame
    if name:
        shape.name = name
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = px(8)
    frame.margin_top = frame.margin_bottom = px(4)
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(vertical, MSO_ANCHOR.MIDDLE)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_card(
    slide,
    box: tuple[float, float, float, float],
    *,
    fill: str = "#FFFDF7",
    transparency: int = 8,
    line: str | None = None,
    name: str | None = None,
):
    """添加可编辑圆角信息卡，透明度使用 PowerPoint 百分比。"""

    left, top, width, height = (px(item) for item in box)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    _set_fill_transparency(shape, transparency)
    if line:
        shape.line.color.rgb = rgb(line)
    else:
        shape.line.fill.background()
    return shape


def _set_fill_transparency(shape, transparency: int) -> None:
    """通过稳定 OOXML alpha 写入 python-pptx 尚未公开的填充透明度。"""

    solid_fill = shape._element.spPr.solidFill
    color_nodes = list(solid_fill)
    if not color_nodes:
        return
    color = color_nodes[0]
    for child in list(color):
        if child.tag.endswith("}alpha"):
            color.remove(child)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(max(0, min(100, 100 - transparency)) * 1000))
    color.append(alpha)


def add_cover_image(slide, image_path: Path, box: tuple[float, float, float, float]) -> None:
    """按 cover 语义裁切图片，避免拉伸并铺满目标区域。"""

    left, top, width, height = box
    with Image.open(image_path) as image:
        source_ratio = image.width / image.height
    target_ratio = width / height
    if source_ratio > target_ratio:
        visible = target_ratio / source_ratio
        crop = (1 - visible) / 2
        picture = slide.shapes.add_picture(str(image_path), px(left), px(top), px(width), px(height))
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        visible = source_ratio / target_ratio
        crop = (1 - visible) / 2
        picture = slide.shapes.add_picture(str(image_path), px(left), px(top), px(width), px(height))
        picture.crop_top = crop
        picture.crop_bottom = crop


def set_internal_link(shape, target_slide) -> None:
    """把可点击形状链接到同一演示文稿中的目标页。"""

    shape.click_action.target_slide = target_slide


def set_external_link(shape, target: str) -> None:
    """设置 URL 或同目录文件链接。"""

    shape.click_action.hyperlink.address = target


@dataclass(frozen=True)
class PreviewResult:
    """记录可视化预览产物；缺少 LibreOffice 时明确返回空结果。"""

    first_page: Path | None
    montage: Path | None


def render_preview(pptx_path: Path, output_dir: Path) -> PreviewResult:
    """使用 LibreOffice 渲染 PDF，再生成首页和整册蒙太奇。"""

    office = shutil.which("libreoffice") or shutil.which("soffice")
    if not office:
        return PreviewResult(None, None)
    output_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [office, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_path)],
        check=True,
        capture_output=True,
        text=True,
    )
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    converter = shutil.which("pdftoppm")
    if not converter or not pdf_path.is_file():
        return PreviewResult(None, None)
    with tempfile.TemporaryDirectory(prefix="selplat_ppt_preview_") as temp_dir:
        prefix = Path(temp_dir) / "page"
        subprocess.run([converter, "-png", "-r", "72", str(pdf_path), str(prefix)], check=True)
        pages = sorted(Path(temp_dir).glob("page-*.png"))
        if not pages:
            return PreviewResult(None, None)
        first_page = output_dir / "首页.png"
        shutil.copy2(pages[0], first_page)
        montage = output_dir / "整册蒙太奇.webp"
        _write_montage(pages, montage)
        return PreviewResult(first_page, montage)


def _write_montage(page_paths: Iterable[Path], target: Path) -> None:
    """把全部页面缩略图排成稳定网格，供人工快速检查。"""

    thumbnails: list[Image.Image] = []
    for path in page_paths:
        with Image.open(path) as image:
            copy = image.convert("RGB")
            copy.thumbnail((320, 180))
            thumbnails.append(copy)
    columns = 4
    rows = max(1, (len(thumbnails) + columns - 1) // columns)
    canvas = Image.new("RGB", (columns * 330, rows * 190), "white")
    draw = ImageDraw.Draw(canvas)
    for index, image in enumerate(thumbnails):
        x = (index % columns) * 330 + 5
        y = (index // columns) * 190 + 5
        canvas.paste(image, (x, y))
        draw.text((x, y + 180), str(index + 1), fill="#333333")
    canvas.save(target, "WEBP", quality=82)
